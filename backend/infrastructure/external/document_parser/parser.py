# input: python-pptx, pdfplumber, python-docx 三方库
# output: FileDocumentParser 文档解析实现
# owner: wanhua.gu
# pos: 基础设施层 - 文档解析适配器 (PPT/PDF/Word → DocumentSummary)；一旦我被更新，务必更新我的开头注释以及所属文件夹的md
"""Concrete document parser: PPT, PDF, Word → DocumentSummary."""
from __future__ import annotations

import io
import re
from pathlib import Path

from domain.defense_prep.value_objects import DocumentSummary, Section

_NUMBER_PATTERN = re.compile(
    r"\d+(?:\.\d+)?%" r"|[$¥€]\s?\d[\d,]*(?:\.\d+)?" r"|\d[\d,]*(?:\.\d+)?\s?[万亿kKmMbB]"
)


class FileDocumentParser:
    """Parse uploaded documents (PPT/PDF/Word/txt) into *DocumentSummary*."""

    async def parse(self, content: bytes, filename: str) -> DocumentSummary:
        ext = Path(filename).suffix.lower()
        if ext == ".pptx":
            return self._parse_pptx(content, filename)
        elif ext == ".pdf":
            return self._parse_pdf(content, filename)
        elif ext == ".docx":
            return self._parse_docx(content, filename)
        elif ext in (".txt", ".md"):
            return self._parse_text(content, filename)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _extract_key_data(self, text: str) -> list[str]:
        """Return deduplicated numeric data points found in *text*."""
        return list(dict.fromkeys(_NUMBER_PATTERN.findall(text)))

    def _parse_pptx(self, content: bytes, filename: str) -> DocumentSummary:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        sections: list[Section] = []
        all_text_parts: list[str] = []
        for slide in prs.slides:
            title = ""
            bullets: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if not title and para.level == 0:
                        title = text
                    else:
                        bullets.append(text)
                    all_text_parts.append(text)
            if title or bullets:
                sections.append(Section(title=title, bullet_points=bullets))
        raw_text = "\n".join(all_text_parts)
        return DocumentSummary(
            title=Path(filename).stem,
            sections=sections,
            key_data=self._extract_key_data(raw_text),
            raw_text=raw_text,
        )

    def _parse_pdf(self, content: bytes, filename: str) -> DocumentSummary:
        import pdfplumber

        all_text_parts: list[str] = []
        sections: list[Section] = []
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            for page in pdf.pages:
                text = page.extract_text() or ""
                all_text_parts.append(text)
                lines = [l.strip() for l in text.split("\n") if l.strip()]
                if lines:
                    sections.append(Section(title=lines[0], bullet_points=lines[1:]))
        raw_text = "\n".join(all_text_parts)
        return DocumentSummary(
            title=Path(filename).stem,
            sections=sections,
            key_data=self._extract_key_data(raw_text),
            raw_text=raw_text,
        )

    def _parse_docx(self, content: bytes, filename: str) -> DocumentSummary:
        from docx import Document

        doc = Document(io.BytesIO(content))
        sections: list[Section] = []
        current_section: Section | None = None
        all_text_parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            all_text_parts.append(text)
            if para.style and para.style.name.startswith("Heading"):
                if current_section:
                    sections.append(current_section)
                current_section = Section(title=text)
            elif current_section:
                current_section.bullet_points.append(text)
            else:
                current_section = Section(title=text)
        if current_section:
            sections.append(current_section)
        raw_text = "\n".join(all_text_parts)
        return DocumentSummary(
            title=Path(filename).stem,
            sections=sections,
            key_data=self._extract_key_data(raw_text),
            raw_text=raw_text,
        )

    def _parse_text(self, content: bytes, filename: str) -> DocumentSummary:
        raw_text = content.decode("utf-8", errors="replace")
        return DocumentSummary(
            title=Path(filename).stem,
            sections=[],
            key_data=self._extract_key_data(raw_text),
            raw_text=raw_text,
        )
